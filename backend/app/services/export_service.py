from __future__ import annotations

import html
import json
import shutil
import uuid
import zipfile
from pathlib import Path

import httpx
from sqlalchemy.orm import Session

from app.core.config import get_settings
from app.core.database import SessionLocal
from app.models import ArtifactVersion, ExportJob, LessonArtifact, ProjectImage
from app.services.theme_service import theme_catalog


EXPORT_REQUIRED_TYPES = {
    "teacher": {"lesson_plan", "slide_deck", "speaker_notes", "exercise_set"},
    "student": {"slide_deck", "exercise_set"},
    "pptx": {"slide_deck"},
}


def _plain_markdown(markdown: str) -> list[str]:
    rows: list[str] = []
    in_comment = False
    for raw in markdown.splitlines():
        clean = raw.strip()
        if clean.startswith("<!--"):
            in_comment = True
        if not in_comment and clean:
            clean = clean.lstrip("#").strip()
            clean = clean.replace("**", "").replace("`", "")
            if clean.startswith(">"):
                clean = clean[1:].strip()
            rows.append(clean)
        if clean.endswith("-->"):
            in_comment = False
    return rows


def _write_pptx(
    path: Path,
    slide_deck: dict,
    image_paths: dict[str, Path] | None = None,
) -> None:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    blank = presentation.slide_layouts[6]
    image_paths = image_paths or {}
    for index, item in enumerate(slide_deck.get("slides", []), 1):
        slide = presentation.slides.add_slide(blank)
        background = slide.background.fill
        background.solid()
        background.fore_color.rgb = RGBColor(246, 248, 252)

        accent = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.18), Inches(7.5))
        accent.fill.solid()
        accent.fill.fore_color.rgb = RGBColor(82, 105, 236)
        accent.line.fill.background()

        placements = item.get("images", [])
        for placement in placements:
            if placement.get("position") != "background":
                continue
            image_path = image_paths.get(placement.get("image_id", ""))
            if image_path and image_path.is_file():
                slide.shapes.add_picture(str(image_path), 0, 0, presentation.slide_width, presentation.slide_height)

        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.65), Inches(11.7), Inches(0.9))
        title_frame = title_box.text_frame
        title_frame.text = item.get("title") or f"第 {index} 页"
        title_frame.paragraphs[0].font.size = Pt(30)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(31, 41, 55)

        body_box = slide.shapes.add_textbox(Inches(0.9), Inches(1.75), Inches(11.4), Inches(4.9))
        body_frame = body_box.text_frame
        body_frame.word_wrap = True
        lines = _plain_markdown(item.get("markdown", ""))
        if lines and lines[0] == item.get("title"):
            lines = lines[1:]
        for line_index, line in enumerate(lines[:16]):
            paragraph = body_frame.paragraphs[0] if line_index == 0 else body_frame.add_paragraph()
            paragraph.text = line.lstrip("- ").strip()
            paragraph.level = 0
            paragraph.font.size = Pt(20 if len(lines) < 9 else 16)
            paragraph.font.color.rgb = RGBColor(55, 65, 81)
            paragraph.space_after = Pt(9)

        number = slide.shapes.add_textbox(Inches(11.9), Inches(6.85), Inches(0.6), Inches(0.3))
        number.text_frame.text = str(index)
        number.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
        number.text_frame.paragraphs[0].font.size = Pt(10)
        number.text_frame.paragraphs[0].font.color.rgb = RGBColor(148, 163, 184)

        for placement in placements:
            if placement.get("position") == "background":
                continue
            image_path = image_paths.get(placement.get("image_id", ""))
            if not image_path or not image_path.is_file():
                continue
            slide.shapes.add_picture(
                str(image_path),
                Inches(13.333 * placement.get("x", 10) / 100),
                Inches(7.5 * placement.get("y", 20) / 100),
                Inches(13.333 * placement.get("width", 40) / 100),
                Inches(7.5 * placement.get("height", 60) / 100),
            )

        note = item.get("speaker_note") or {}
        note_lines = [
            note.get("explanation", ""),
            "课堂提问：" + "；".join(note.get("questions", [])) if note.get("questions") else "",
            "板书：" + note.get("board_notes", "") if note.get("board_notes") else "",
            "过渡：" + note.get("transition", "") if note.get("transition") else "",
        ]
        try:
            slide.notes_slide.notes_text_frame.text = "\n".join(row for row in note_lines if row)
        except (AttributeError, ValueError):
            pass
    presentation.save(path)


def _slidev_notes(slide: dict) -> str:
    note = slide.get("speaker_note") or {}
    rows = [
        note.get("explanation", ""),
        "课堂提问：" + "；".join(note.get("questions", [])) if note.get("questions") else "",
        "参考回答：" + "；".join(note.get("expected_answers", [])) if note.get("expected_answers") else "",
        "板书：" + note.get("board_notes", "") if note.get("board_notes") else "",
        "过渡：" + note.get("transition", "") if note.get("transition") else "",
    ]
    value = "\n".join(row for row in rows if row)
    return f"\n\n<!--\n{value}\n-->" if value else ""


def _placement_html(placement: dict, asset_name: str) -> str:
    position = placement.get("position", "right")
    opacity = placement.get("opacity", 1)
    style = (
        "position:absolute;"
        f"left:{placement.get('x', 10)}%;top:{placement.get('y', 20)}%;"
        f"width:{placement.get('width', 40)}%;height:{placement.get('height', 60)}%;"
        f"object-fit:cover;opacity:{opacity};"
        f"z-index:{0 if position == 'background' else 3};"
        "border-radius:14px;"
    )
    caption = html.escape(placement.get("caption", ""))
    image = f'<img src="/assets/{html.escape(asset_name)}" style="{style}" />'
    if not caption or position == "background":
        return f'<div class="any2ppt-image-layer">{image}</div>'
    caption_style = (
        "position:absolute;"
        f"left:{placement.get('x', 10)}%;"
        f"top:{min(94, placement.get('y', 20) + placement.get('height', 60) + 1)}%;"
        f"width:{placement.get('width', 40)}%;"
        "text-align:center;font-size:12px;opacity:.75;z-index:4;"
    )
    return (
        '<div class="any2ppt-image-layer">'
        + image
        + f'<div style="{caption_style}">{caption}</div>'
        + "</div>"
    )


def _prepare_slidev_job(
    job_dir: Path,
    slide_deck: dict,
    image_records: dict[str, ProjectImage],
) -> None:
    public_assets = job_dir / "public" / "assets"
    public_assets.mkdir(parents=True, exist_ok=True)
    asset_names: dict[str, str] = {}
    for image_id, record in image_records.items():
        suffix = Path(record.storage_path).suffix.lower()
        asset_name = f"{image_id}{suffix}"
        shutil.copy2(record.storage_path, public_assets / asset_name)
        asset_names[image_id] = asset_name

    theme_package = slide_deck.get("theme", "@slidev/theme-default")
    slides = slide_deck.get("slides", [])
    frontmatter_rows = [
        "---",
        f"theme: {json.dumps(theme_package, ensure_ascii=False)}",
        f"title: {json.dumps(slide_deck.get('deck_title', '备课课件'), ensure_ascii=False)}",
        "download: false",
        "canvasWidth: 1280",
        "aspectRatio: 16/9",
    ]
    if slide_deck.get("theme_config"):
        frontmatter_rows.append(
            f"themeConfig: {json.dumps(slide_deck['theme_config'], ensure_ascii=False)}"
        )
    if slides:
        frontmatter_rows.append(
            f"layout: {json.dumps(slides[0].get('layout', 'default'), ensure_ascii=False)}"
        )
    frontmatter_rows.append("---")
    frontmatter = "\n".join(frontmatter_rows)
    pages = []
    for slide in slides:
        placements = slide.get("images", [])
        backgrounds = []
        foregrounds = []
        for placement in placements:
            asset_name = asset_names.get(placement.get("image_id", ""))
            if not asset_name:
                continue
            target = backgrounds if placement.get("position") == "background" else foregrounds
            target.append(_placement_html(placement, asset_name))
        pages.append(
            "\n\n".join(backgrounds)
            + "\n\n"
            + slide.get("markdown", "")
            + "\n\n"
            + "\n\n".join(foregrounds)
            + _slidev_notes(slide)
        )
    document = frontmatter
    if pages:
        document += "\n\n" + pages[0]
    for slide, page in zip(slides[1:], pages[1:]):
        document += (
            "\n\n---\n"
            f"layout: {json.dumps(slide.get('layout', 'default'), ensure_ascii=False)}\n"
            "---\n\n"
            + page
        )
    (job_dir / "slides.md").write_text(document, "utf-8")


def _write_slidev_pptx(
    path: Path,
    slide_deck: dict,
    image_records: dict[str, ProjectImage],
    project_id: str | None = None,
) -> None:
    settings = get_settings()
    if not settings.slidev_renderer_url:
        _write_pptx(
            path,
            slide_deck,
            {image_id: Path(record.storage_path) for image_id, record in image_records.items()},
        )
        return
    allowed = {
        (item["package"], item["version"])
        for item in theme_catalog()
    }
    theme_package = slide_deck.get("theme", "@slidev/theme-default")
    theme_version = slide_deck.get("theme_version", "0.25.0")
    if (theme_package, theme_version) not in allowed:
        raise ValueError("THEME_NOT_ALLOWED")

    job_id = str(uuid.uuid4())
    job_dir = settings.export_dir.resolve().parent / "render_jobs" / job_id
    job_dir.mkdir(parents=True, exist_ok=False)
    try:
        _prepare_slidev_job(job_dir, slide_deck, image_records)
        response = httpx.post(
            f"{settings.slidev_renderer_url.rstrip('/')}/render",
            json={
                "job_id": job_id,
                "project_id": project_id,
                "theme_package": theme_package,
                "theme_version": theme_version,
            },
            timeout=settings.slidev_renderer_timeout_seconds,
        )
        response.raise_for_status()
        rendered = job_dir / "output.pptx"
        if not rendered.is_file():
            raise RuntimeError("SLIDEV_RENDER_OUTPUT_MISSING")
        shutil.move(rendered, path)
    finally:
        shutil.rmtree(job_dir, ignore_errors=True)


def _load_versions(db: Session, job: ExportJob) -> dict[str, ArtifactVersion]:
    selected_ids = list((job.selected_versions or {}).values())
    if selected_ids:
        versions = (
            db.query(ArtifactVersion)
            .join(LessonArtifact, ArtifactVersion.artifact_id == LessonArtifact.id)
            .filter(
                LessonArtifact.project_id == job.project_id,
                ArtifactVersion.id.in_(selected_ids),
            )
            .all()
        )
    else:
        artifacts = db.query(LessonArtifact).filter_by(project_id=job.project_id).all()
        versions = [artifact.versions[-1] for artifact in artifacts if artifact.versions]

    by_type = {version.artifact.type: version for version in versions}
    required = EXPORT_REQUIRED_TYPES[job.package_type]
    missing = sorted(required - by_type.keys())
    if missing:
        raise ValueError(f"EXPORT_ARTIFACTS_MISSING: {', '.join(missing)}")
    return {artifact_type: by_type[artifact_type] for artifact_type in sorted(required)}


def _write_slides(output: Path, slide_deck: dict) -> None:
    slides = slide_deck.get("slides", [])
    frontmatter = (
        "---\n"
        f"theme: {slide_deck.get('theme', 'seriph')}\n"
        f"title: {json.dumps(slide_deck.get('deck_title', 'LessonDeck'), ensure_ascii=False)}\n"
        "download: false\n"
        "---"
    )
    slidev_source = frontmatter + "\n\n" + "\n\n---\n\n".join(
        slide.get("markdown", "") for slide in slides
    )
    (output / "slides.md").write_bytes(slidev_source.encode("utf-8"))

    def render_markdown(source: str) -> str:
        lines = []
        for raw in source.splitlines():
            clean = html.escape(raw)
            if clean.startswith("# "):
                lines.append(f"<h1>{clean[2:]}</h1>")
            elif clean.startswith("## "):
                lines.append(f"<h2>{clean[3:]}</h2>")
            elif clean.startswith("&gt; "):
                lines.append(f"<blockquote>{clean[5:]}</blockquote>")
            elif clean.strip():
                lines.append(f"<p>{clean}</p>")
        return "".join(lines)

    sections = "".join(
        f"<section id='{html.escape(slide.get('slide_id', ''))}'>"
        f"{render_markdown(slide.get('markdown', ''))}</section>"
        for slide in slides
    )
    page = (
        "<!doctype html><meta charset='utf-8'><title>LessonDeck</title>"
        "<style>"
        "body{font-family:sans-serif;background:#eef2f8;margin:0}"
        "section{box-sizing:border-box;width:100vw;height:100vh;"
        "padding:8vh 10vw;background:white;border-bottom:1px solid #ccd3df}"
        "h1{font-size:5vw}h2{font-size:3.5vw}"
        "p,blockquote{font:2vw/1.6 sans-serif}blockquote{border-left:5px solid #49a;padding-left:2vw}"
        "@media print{section{page-break-after:always}}"
        "</style>"
        f"{sections}"
    )
    (output / "slides.html").write_text(page, "utf-8")


def create_export(job_id: str) -> None:
    db = SessionLocal()
    job = db.get(ExportJob, job_id)
    if not job:
        db.close()
        return

    root = get_settings().export_dir.resolve()
    root.mkdir(parents=True, exist_ok=True)
    tmp = root / "_tmp" / job.id
    if tmp.exists():
        shutil.rmtree(tmp, ignore_errors=True)
    tmp.mkdir(parents=True, exist_ok=False)

    try:
        job.status = "running"
        db.commit()

        versions = _load_versions(db, job)
        latest = {artifact_type: version.content for artifact_type, version in versions.items()}
        if job.package_type == "pptx":
            pptx_path = root / f"{job.project_id}_{job.id}.pptx"
            tmp_pptx = tmp / pptx_path.name
            image_ids = {
                placement.get("image_id")
                for slide in latest["slide_deck"].get("slides", [])
                for placement in slide.get("images", [])
                if placement.get("image_id")
            }
            image_rows = (
                db.query(ProjectImage)
                .filter(
                    ProjectImage.project_id == job.project_id,
                    ProjectImage.id.in_(image_ids),
                )
                .all()
                if image_ids
                else []
            )
            image_records = {item.id: item for item in image_rows}
            if missing_images := sorted(image_ids - image_records.keys()):
                raise ValueError(f"EXPORT_IMAGES_MISSING: {', '.join(missing_images)}")
            _write_slidev_pptx(tmp_pptx, latest["slide_deck"], image_records, job.project_id)
            tmp_pptx.replace(pptx_path)
            job.selected_versions = {"slide_deck": versions["slide_deck"].id}
            job.file_path = str(pptx_path)
            job.status = "succeeded"
            db.commit()
            return

        output = tmp / "package"
        output.mkdir()
        (output / "README.txt").write_text(
            "LessonDeck 备课导出包\nslides.md 为 Slidev 兼容源码；slides.html 可离线预览或打印为 PDF。\n内容由教师最终确认后使用。",
            "utf-8",
        )
        (output / "版本清单.json").write_text(
            json.dumps(
                {
                    artifact_type: {
                        "version_id": version.id,
                        "version_no": version.version_no,
                        "change_type": version.change_type,
                    }
                    for artifact_type, version in versions.items()
                },
                ensure_ascii=False,
                indent=2,
            ),
            "utf-8",
        )

        if "slide_deck" in latest:
            _write_slides(output, latest["slide_deck"])

        if job.package_type == "teacher":
            names = {
                "lesson_plan": "教学设计.json",
                "speaker_notes": "逐页讲稿.json",
                "exercise_set": "教师版练习.json",
            }
            for artifact_type, filename in names.items():
                (output / filename).write_text(
                    json.dumps(latest[artifact_type], ensure_ascii=False, indent=2),
                    "utf-8",
                )
            citations = [
                citation
                for version in versions.values()
                for citation in (version.citations or [])
            ]
            (output / "引用清单.json").write_text(
                json.dumps(citations, ensure_ascii=False, indent=2),
                "utf-8",
            )
        else:
            student = json.loads(json.dumps(latest["exercise_set"], ensure_ascii=False))
            for item in student.get("exercises", []):
                item.pop("answer", None)
                item.pop("explanation", None)
            (output / "学生练习.json").write_text(
                json.dumps(student, ensure_ascii=False, indent=2),
                "utf-8",
            )

        zip_path = root / f"{job.project_id}_{job.package_type}_{job.id}.zip"
        tmp_zip = tmp / zip_path.name
        with zipfile.ZipFile(tmp_zip, "w", zipfile.ZIP_DEFLATED) as archive:
            for file in output.iterdir():
                archive.write(file, file.name)
        tmp_zip.replace(zip_path)

        job.selected_versions = {artifact_type: version.id for artifact_type, version in versions.items()}
        job.file_path = str(zip_path)
        job.status = "succeeded"
        db.commit()
    except Exception as exc:
        job.status = "failed"
        job.error_message = str(exc)[:500]
        db.commit()
    finally:
        shutil.rmtree(tmp, ignore_errors=True)
        db.close()
